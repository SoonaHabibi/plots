# -*- coding: utf-8 -*-
"""
Created on Mon Apr  4 15:22:34 2022
@author: cyu modified by Soona

For CRT event selection review. 
Plots each calibration point's set of chosen winter and spring (if applicable) event timeseries with three subplots: 

    
Ex: Plot of Brownlee Winter:     
    Subplot 1: Precipitation (10 timeseries in red for chosen years. All other years plotted in gray.)
    Subplot 2: Tmax (10 timeseries in red for chosen years. All other years plotted in gray.)
    Subplot 3: Streamflow (10 timeseries in red for chosen years. All other years plotted in gray.)
    
Note: Winter and Spring season dates according to flow season. (also here: P:\2021\CRT_Task3\Data\CalibPointSeasonality.csv)
    Winter: Oct 1- Apr1 (flow season)
    Spring: Apr1 - Oct1 (flow season)
    
!!! SWITCH BETWEEN NESTED AND INDIVIDUAL USING FLAG "flag_events"

Finally you can write all plots into a word filr


"""
import sys 


sys.path.append(r"O:\Python\WEST_Python_FunctionAug2019")
sys.path.append("P:\2021\CRT_Task3")
import BasicFunction_py3 as BF   
import geopandas as gpd
import numpy as np
import os
import pandas as pd
from osgeo import ogr
from pathlib import Path
import matplotlib.pyplot as plt
import glob
from matplotlib.lines import Line2D
import datetime
from PIL import Image

sys.path.append(r"P:\2021\CRT_Task3\Code");import CDA_config #configuration file


#%%Directories
flag_events = "indiv" #options: "indiv" or "nested"
plotTogether = False
indir_flow = CDA_config.indir_flow
flow_rain_link_file =CDA_config.flow_rain_link_file # r'D:\CRT_Task3\Data\flow_rain_calibpoint_linkage_TESTONLY.csv'

cp_tmax_dir= CDA_config.cp_avTmax
cp_rain_dir= CDA_config.cp_avRain
cp_swe_dir = CDA_config.cp_avSWE
event_dir = "//westfolsom/projects/2021/CRT_Task3/Data/SelectedEvents.csv"# These are the final selected events for nested basins
event_indiv= "//westfolsom/projects/2021/CRT_Task3/Data/IndividualSelection.csv"


season_def_file = CDA_config.season_def_file


#Read in selected events
outdir_nested = r"\\westfolsom\projects\2021\CRT_Task3\MaxEventsYearSelection\SelectedEventsTempPrecipFlowVariability/"
outdir_indiv =  r"\\westfolsom\projects\2021\CRT_Task3\MaxEventsYearSelection\SelectedEventsTempPrecipFlowVariability\IndivBasinSelectionOnly/"
outdir_both = r"\\westfolsom\projects\2021\CRT_Task3\MaxEventsYearSelection\SelectedEventsTempPrecipFlowVariability\bothIndAndNest/"

if flag_events== "nested":
    events_df = pd.read_csv(event_dir)
    outdir = outdir_nested
    plot_num=['(n.1)','(n.2)','(n.3)','(n.4)']
    
elif flag_events == "indiv":
    events_df = pd.read_csv(event_indiv)
    events_df= events_df.iloc[:10]
    outdir = outdir_indiv
    plot_num=['(i.1)','(i.2)','(i.3)','(i.4)']


if not os.path.exists(outdir): os.mkdir(outdir)

#%%Create plots

#reformat the column names so it matches cp names in other files
col_names = events_df.columns.tolist()
col_names = [col.lower().replace("spaulding", "spalding") for col in col_names ]
events_df.columns = col_names

seasons = pd.read_csv(season_def_file, index_col = 0)
calibration_points = seasons.index.tolist()

flow_rain_link = pd.read_csv(flow_rain_link_file, index_col = 0) 
cp_flow = dict(zip(flow_rain_link.index.tolist(), flow_rain_link['unimpaired flow with natural lakes'].astype(str).tolist())) #create a dict with calib point: flow name

#Read in flow data
flow_cfs= pd.read_csv(indir_flow,skiprows=[1,2,3,4,5,6], index_col = 1)
flow_cfs.rename(columns = {'A': 'Record_Number'}, inplace = True)
flow_cfs.index = pd.to_datetime (flow_cfs.index)
flow_cfs['wateryear']= flow_cfs.index.map(lambda x: x.year if (x.month < 10) else (x.year+1))
flow_cfs['dayofyear']= flow_cfs.index.map (lambda x: x.timetuple().tm_yday)
flow_cfs['dayofwateryear']= flow_cfs.index.map(lambda x: x.timetuple().tm_yday- 273 if (x.year%4 !=0 and x.timetuple().tm_yday >= 274) else ( x.timetuple().tm_yday+92  if (x.year%4 !=0 and x.timetuple().tm_yday < 274) else (x.timetuple().tm_yday- 274 if x.timetuple().tm_yday>274 else x.timetuple().tm_yday + 92) ))





#Start loop through each calibration point
for cp in calibration_points: 
    print(cp)
    num_seasons = seasons.at[cp, "NumSeasons"] #Either 1 or two seasons

    #Get list of events selected
    if num_seasons ==2: 
        events_winter = events_df[cp + "_winter"].dropna().astype(int).tolist()
        events_winter.sort()
        winter_start= seasons.at[cp, 'WinterStart']
        winter_end = seasons.at[cp, 'WinterEnd']
        
    
    events_spring = events_df[cp].dropna().astype(int).tolist()
    events_spring.sort()
    spring_start = seasons.at[cp, 'SpringStart']
    spring_end = seasons.at[cp, "SpringEnd"]
    
    #Read in rain, tmax, and flow for the calibration point
    rainfile= cp_rain_dir  + cp.replace(' ', '_')+'_mm.csv'#Basin averaged rain for the calibration point area
    tmaxfile = cp_tmax_dir + cp.replace(" ", "_")+ "_F.csv"
    swefile = cp_swe_dir +cp.replace(' ', '_')+'_mm.csv'
    
    rain= pd.read_csv(rainfile, index_col = "end_date", parse_dates = True)
    rain['wateryear']= rain.index.map(lambda x: x.year if (x.month < 10) else (x.year+1))
    rain['dayofyear']= rain.index.map (lambda x: x.timetuple().tm_yday)
    rain['dayofwateryear']= rain.index.map(lambda x: x.timetuple().tm_yday- 273 if (x.year%4 !=0 and x.timetuple().tm_yday >= 274) else ( x.timetuple().tm_yday+92  if (x.year%4 !=0 and x.timetuple().tm_yday < 274) else (x.timetuple().tm_yday- 274 if x.timetuple().tm_yday>274 else x.timetuple().tm_yday + 92) ))
    rain_dated = rain.reset_index().set_index("dayofwateryear")
    
    tmax = pd.read_csv(tmaxfile, index_col = "end_date", parse_dates = True)
    tmax['wateryear']= tmax.index.map(lambda x: x.year if (x.month < 10) else (x.year+1))
    tmax['dayofyear']= tmax.index.map (lambda x: x.timetuple().tm_yday)
    tmax['dayofwateryear']= tmax.index.map(lambda x: x.timetuple().tm_yday- 273 if (x.year%4 !=0 and x.timetuple().tm_yday >= 274) else ( x.timetuple().tm_yday+92  if (x.year%4 !=0 and x.timetuple().tm_yday < 274) else (x.timetuple().tm_yday- 274 if x.timetuple().tm_yday>274 else x.timetuple().tm_yday + 92) ))
    tmax_dated=tmax.reset_index().set_index ('dayofwateryear')

    swe= pd.read_csv(swefile, index_col = "END_DATE", parse_dates = True)
    swe.index.name = "end_date"
    swe['wateryear']= swe.index.map(lambda x: x.year if (x.month < 10) else (x.year+1))
    swe['dayofyear']= swe.index.map (lambda x: x.timetuple().tm_yday)
    swe['dayofwateryear']= swe.index.map(lambda x: x.timetuple().tm_yday- 273 if (x.year%4 !=0 and x.timetuple().tm_yday >= 274) else ( x.timetuple().tm_yday+92  if (x.year%4 !=0 and x.timetuple().tm_yday < 274) else (x.timetuple().tm_yday- 274 if x.timetuple().tm_yday>274 else x.timetuple().tm_yday + 92) ))
    swe_dated = swe.reset_index().set_index("dayofwateryear")
 
    flow_cp = flow_cfs[[cp_flow[cp], 'wateryear', 'dayofyear', 'dayofwateryear']]
    flow_cp.index.name = "end_date"
    flow_cp_dated= flow_cp.reset_index().set_index ('dayofwateryear') 

    #Find the years that there is data for rain, tmax, and flow
    yearMin =max(min(flow_cfs.wateryear.min(),rain.wateryear.min()), tmax.wateryear.min())
    yearMax =  min(min(flow_cfs.wateryear.max(),rain.wateryear.max()), tmax.wateryear.max())
    years = range(yearMin, yearMax + 1)
    
    
    numSubplots=4

    #Plot 
    if plotTogether: 
        fig, ax = plt.subplots(numSubplots, num_seasons, figsize = (min(20*num_seasons, 25),30))
        fig.tight_layout(h_pad=2.5, w_pad=2.5)
        # font size
        plt.rc('font', size=21) 
        # plt.subplots_adjust(top=0.9)
    else: 
        fig, ax = plt.subplots(numSubplots, 1, figsize = (20,20))
        fig_w, ax_w = plt.subplots(numSubplots, 1, figsize = (20,20))
        fig.tight_layout(h_pad=2.5)
        fig_w.tight_layout(h_pad=2.5)
        plt.rc('font', size=21) 
        # plt.subplots_adjust(top=0.9)

    spring_years= []
    winter_years = []
    for y in years: #Loop through each year
        
        if num_seasons==1: 
            if y in events_spring:
                spring_years.append(y)
                linecolor = 'r'
                z = 2
            else:
                linecolor = '0.75'
                z = 1
            ax[0].plot(rain_dated.loc[(rain_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (rain_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
            ax[0].text(-.07, -0.025, plot_num[0], fontsize=22,fontweight='bold', color='royalblue', transform=ax[0].transAxes)
            ax[1].plot(tmax_dated.loc[(tmax_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (tmax_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
            ax[1].text(-.07, -0.025, plot_num[1], fontsize=22,fontweight='bold', color='royalblue', transform=ax[1].transAxes)
            ax[2].plot(flow_cp_dated.loc[(flow_cp_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (flow_cp_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y")), cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
            ax[2].text(-.07, -0.025, plot_num[2], fontsize=22,fontweight='bold', color='royalblue', transform=ax[2].transAxes)
            ax[3].plot(swe_dated.loc[(swe_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (swe_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
            ax[3].text(-.07, -0.025, plot_num[3], fontsize=22,fontweight='bold', color='royalblue', transform=ax[3].transAxes)
            for i in range(numSubplots):
                ax[i].axvline(x = 183, linestyle = '--', color = 'black') #Add a line at day April 1 (non-leap year)
                ax[i].axvline(x = 244, linestyle = '--', color = 'black') #Jun 1
                ax[i].axvline(x = 305, linestyle = "--", color = 'black' ) #Aug 1
        elif num_seasons ==2: 
            #Spring

            if y in events_spring:
                spring_years.append(y)
                linecolor= 'r'
                z = 2
            else:
                linecolor = '0.75'
                z = 1
            
            
            if plotTogether: 
                i = 1
                ax[0, i].plot(rain_dated.loc[(rain_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (rain_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
                ax[1, i].plot(tmax_dated.loc[(tmax_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (tmax_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax[2, i].plot(flow_cp_dated.loc[(flow_cp_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (flow_cp_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y")), cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
                ax[3, i].plot(swe_dated.loc[(swe_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (swe_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
            else: 
                ax[0].plot(rain_dated.loc[(rain_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (rain_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
                ax[0].text(-.07, -0.025, plot_num[0], fontsize=22,fontweight='bold', color='royalblue', transform=ax[0].transAxes)
                ax[1].plot(tmax_dated.loc[(tmax_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (tmax_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax[1].text(-.07, -0.025, plot_num[1], fontsize=22,fontweight='bold', color='royalblue', transform=ax[1].transAxes)
                ax[2].plot(flow_cp_dated.loc[(flow_cp_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (flow_cp_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y")), cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
                ax[2].text(-.07, -0.025, plot_num[2], fontsize=22,fontweight='bold', color='royalblue', transform=ax[2].transAxes)
                ax[3].plot(swe_dated.loc[(swe_dated.end_date>= datetime.datetime.strptime(spring_start +"-"+ str(y), "%d-%b-%Y"))& (swe_dated.end_date< datetime.datetime.strptime(spring_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax[3].text(-.07, -0.025, plot_num[3], fontsize=22,fontweight='bold', color='royalblue', transform=ax[3].transAxes)
                
            # ax[2, i].plot(flow_cp_dated.loc[flow_cp_dated.wateryear == y, cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
            
            #winter
            
            if y in events_winter:
                winter_years.append(y)
                linecolor= 'r'
                z = 2
            else:
                linecolor = '0.75'
                z = 1
            # ax[0, i].plot(rain_dated.loc[rain_dated.wateryear== y]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
            # ax[1, i].plot(tmax_dated.loc[tmax_dated.wateryear == y]['value'], color = linecolor, alpha = 0.5, zorder = z)
            # ax[2, i].plot(flow_cp_dated.loc[flow_cp_dated.wateryear == y, cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
            if plotTogether: 
                i = 0
                ax[0, i].plot(rain_dated.loc[(rain_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (rain_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
                ax[1, i].plot(tmax_dated.loc[(tmax_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (tmax_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax[2, i].plot(flow_cp_dated.loc[(flow_cp_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (flow_cp_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y")), cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
                ax[3, i].plot(swe_dated.loc[(swe_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (swe_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
            else: 
                ax_w[0].plot(rain_dated.loc[(rain_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (rain_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)#swe
                ax_w[0].text(-.07, -0.025, plot_num[0], fontsize=22,fontweight='bold', color='royalblue', transform=ax_w[0].transAxes)
                ax_w[1].plot(tmax_dated.loc[(tmax_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (tmax_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax_w[1].text(-.07, -0.025, plot_num[1], fontsize=22,fontweight='bold', color='royalblue', transform=ax_w[1].transAxes)
                ax_w[2].plot(flow_cp_dated.loc[(flow_cp_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (flow_cp_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y")), cp_flow[cp]], color = linecolor, alpha = 0.5, zorder = z) #streamflow
                ax_w[2].text(-.07, -0.025, plot_num[2], fontsize=22,fontweight='bold', color='royalblue', transform=ax_w[2].transAxes)
                ax_w[3].plot(swe_dated.loc[(swe_dated.end_date>= datetime.datetime.strptime(winter_start +"-"+ str(y-1), "%d-%b-%Y"))& (swe_dated.end_date< datetime.datetime.strptime(winter_end +"-"+ str(y), "%d-%b-%Y"))]['value'], color = linecolor, alpha = 0.5, zorder = z)
                ax_w[3].text(-.07, -0.025, plot_num[3], fontsize=22,fontweight='bold', color='royalblue', transform=ax_w[3].transAxes)
                
            
    for i in range(numSubplots):
        
        if plotTogether: 
            if num_seasons ==2: 
                ax[i,1].axvline(x = 183, linestyle = '--', color = 'black') #Add a line at day April 1 (non-leap year)
                ax[i,1].axvline(x = 244, linestyle = '--', color = 'black') #Jun 1
                ax[i,1].axvline(x =305, linestyle = "--", color = 'black' ) #Aug 1
                
                ax[i,1].annotate("Apr 1",xy=(183+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i,1].annotate("Jun 1",xy=(244+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i,1].annotate("Aug 1",xy=(305+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
    
                
                ax[i,0].axvline(x = 1, linestyle = '--', color = 'black') #Add a line at day Oct1 (non-leap year)
                ax[i,0].axvline(x = 62, linestyle = '--', color = 'black') #Dec1
                ax[i,0].axvline(x =124 , linestyle = "--", color = 'black' ) #Feb1
                
                ax[i,0].annotate("Oct 1",xy=(1+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i,0].annotate("Dec 1",xy=(62+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i,0].annotate("Feb 1",xy=(124+1,ax[i,0].get_ylim()[-1]*2/3 ), xycoords='data')
            else: 
                ax[i].axvline(x = 183, linestyle = '--', color = 'black') #Add a line at day April 1 (non-leap year)
                ax[i].axvline(x = 244, linestyle = '--', color = 'black') #Jun 1
                ax[i].axvline(x =305, linestyle = "--", color = 'black' ) #Aug 1
                
                ax[i].annotate("Apr 1",xy=(183+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i].annotate("Jun 1",xy=(244+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
                ax[i].annotate("Aug 1",xy=(305+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
                
        else: 
            if num_seasons ==2: 
                #Annotate the winter plot
         
                ax_w[i].axvline(x = 1, linestyle = '--', color = 'black') #Add a line at day Oct1 (non-leap year)
                ax_w[i].axvline(x = 62, linestyle = '--', color = 'black') #Dec1
                ax_w[i].axvline(x =124 , linestyle = "--", color = 'black' ) #Feb1
                
                ax_w[i].annotate("Oct 1",xy=(1+1,ax_w[i].get_ylim()[-1]*2/3 ), xycoords='data')
                ax_w[i].annotate("Dec 1",xy=(62+1,ax_w[i].get_ylim()[-1]*2/3 ), xycoords='data')
                ax_w[i].annotate("Feb 1",xy=(124+1,ax_w[i].get_ylim()[-1]*2/3 ), xycoords='data')
            
            ax[i].axvline(x = 183, linestyle = '--', color = 'black') #Add a line at day April 1 (non-leap year)
            ax[i].axvline(x = 244, linestyle = '--', color = 'black') #Jun 1
            ax[i].axvline(x =305, linestyle = "--", color = 'black' ) #Aug 1
            
            ax[i].annotate("Apr 1",xy=(183+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
            ax[i].annotate("Jun 1",xy=(244+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
            ax[i].annotate("Aug 1",xy=(305+1,ax[i].get_ylim()[-1]*2/3 ), xycoords='data')
    
                
            


            
             
    custom_lines = [Line2D([0], [0], color='r', lw=4),
                    Line2D([0], [0], color='0.75', lw=4)]
    #Add axis lables/titles
    if plotTogether: 
        if num_seasons ==1: 
            if spring_years != events_spring: 
                raise ValueError("Spring Evs do not match for " + cp)
            ax[0].set_title("Precipitation", fontweight='bold');
            ax[0].set_ylabel("Precipitation (mm)");
            ax[0].set_xlabel ("Day of water year")
            ax[0].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax[1].set_title("Maximum Daily Temperature", fontweight='bold'); 
            ax[1].set_ylabel("Temperature (degrees F)");
            ax[1].set_xlabel ("Day of water year")
            ax[1].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax[2].set_title("Flow", fontweight='bold')
            ax[2].set_ylabel("Flow(cfs)"); 
            ax[2].set_xlabel ("Day of water year")
            ax[2].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax[3].set_title("SWE", fontweight='bold'); 
            ax[3].set_ylabel("SWE (mm)"); 
            ax[3].set_xlabel ("Day of water year")
            ax[3].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            
            # #Set xlim to days 183 through 366
            # ax[0].set_xlim(183,366) #Roughly April 1 through Oct 1 ("roughly" bc leap years have different day numbers)
            # ax[1].set_xlim(183,366)
            # ax[2].set_xlim(183,366)
        elif num_seasons ==2: 
            season= ["Winter", "Spring" ]
            if spring_years != events_spring: 
                raise ValueError("Spring Evs do not match for " + cp)
            if winter_years != events_winter: 
                raise ValueError("Winter Evs do not match for " + cp)
            
            for i in range(num_seasons): 
                ax[0,i].set_title(cp.title() + ", " + season[i]+"\n Precipitation", fontweight='bold');
                ax[0,i].set_ylabel("Precipitation (mm)");
                ax[0,i].set_xlabel ("Day of water year")
                ax[0,i].legend(custom_lines, ["Selected Years", "All Other Years"])
    
                ax[1,i].set_title("Maximum Daily Temperature", fontweight='bold'); 
                ax[1,i].set_ylabel("Temperature (degrees F)");
                ax[1,i].set_xlabel ("Day of water year")
                ax[1,i].legend(custom_lines, ["Selected Years", "All Other Years"])
    
                ax[2,i].set_title("Flow", fontweight='bold'); 
                ax[2,i].set_ylabel("Flow(cfs)"); 
                ax[2,i].set_xlabel ("Day of water year")
                ax[2,i].legend(custom_lines, ["Selected Years", "All Other Years"])
                
                ax[3,i].set_title("SWE", fontweight='bold'); 
                ax[3,i].set_ylabel("SWE (mm)"); 
                ax[3,i].set_xlabel ("Day of water year")
                ax[3,i].legend(custom_lines, ["Selected Years", "All Other Years"])
                
                # if i == 0: 
                #     xmin = 1
                #     xmax = 183
                # if i ==1: 
                #     xmin = 183
                #     xmax = 366
                # ax[0,i].set_xlim(xmin, xmax)
                # ax[1,i].set_xlim(xmin, xmax)
                # ax[2,i].set_xlim(xmin, xmax)
                
    else:
        if spring_years != events_spring: 
            raise ValueError("Spring Evs do not match for " + cp)
        ax[0].set_title("Precipitation", fontweight='bold');
        ax[0].set_ylabel("Precipitation (mm)");
        ax[0].set_xlabel ("Day of water year")
        ax[0].legend(custom_lines, ["Selected Years", "All Other Years"])
        
        ax[1].set_title("Maximum Daily Temperature", fontweight='bold'); 
        ax[1].set_ylabel("Temperature (degrees F)");
        ax[1].set_xlabel ("Day of water year")
        ax[1].legend(custom_lines, ["Selected Years", "All Other Years"])
        
        ax[2].set_title("Flow", fontweight='bold'); 
        ax[2].set_ylabel("Flow(cfs)"); 
        ax[2].set_xlabel ("Day of water year")
        ax[2].legend(custom_lines, ["Selected Years", "All Other Years"])
        
        ax[3].set_title("SWE", fontweight='bold'); 
        ax[3].set_ylabel("SWE (mm)"); 
        ax[3].set_xlabel ("Day of water year")
        ax[3].legend(custom_lines, ["Selected Years", "All Other Years"])
        # fig.suptitle(cp.title() + " Spring Season", fontweight='bold')
        # plt.subplots_adjust(top=0.96)
        
        
        if num_seasons == 2: 
            if winter_years != events_winter: 
                raise ValueError("Winter Evs do not match for " + cp)
            ax_w[0].set_title("Precipitation", fontweight='bold');
            ax_w[0].set_ylabel("Precipitation (mm)");
            ax_w[0].set_xlabel ("Day of water year")
            ax_w[0].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax_w[1].set_title("Maximum Daily Temperature", fontweight='bold'); 
            ax_w[1].set_ylabel("Temperature (degrees F)");
            ax_w[1].set_xlabel ("Day of water year")
            ax_w[1].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax_w[2].set_title("Flow", fontweight='bold'); 
            ax_w[2].set_ylabel("Flow(cfs)"); 
            ax_w[2].set_xlabel ("Day of water year")
            ax_w[2].legend(custom_lines, ["Selected Years", "All Other Years"])
            
            ax_w[3].set_title("SWE", fontweight='bold'); 
            ax_w[3].set_ylabel("SWE (mm)"); 
            ax_w[3].set_xlabel ("Day of water year")
            ax_w[3].legend(custom_lines, ["Selected Years", "All Other Years"])
            # fig_w.suptitle(cp.title() + " Winter Season", fontweight='bold')
            # plt.subplots_adjust(top=0.95)
    
    
    
            
        
        
    if plotTogether: 
        plt.subplots_adjust(hspace=0.4)
        fig.suptitle(cp.title() , y = 0.95, fontsize = 20, fontweight='bold')
        fig.savefig(outdir + cp.replace(" ", "_") + "_spring.png",bbox_inches='tight')
        
    else: 
        if num_seasons==2: 
            # plt.subplots_adjust(hspace=0.2)
            fig_w.savefig(outdir + cp.replace(" ", "_") + "_winter.png",bbox_inches='tight')
        # plt.subplots_adjust(hspace=0.2)
        fig.savefig(outdir + cp.replace(" ", "_") + "_spring.png",bbox_inches='tight')



#%%%

from pptx import Presentation
from pptx.util import Inches
from pandas.plotting import table 
import seaborn as sns

if plotTogether: 

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]
    print("Generating Summary PPTX ")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Selected Events Precip/Tmax/Flow"
    subtitle.text = "Winter Season: Plotted from Oct 1- Mar 31\nSpring Season: Plotted from Apr 1- Sep 30 (To show full flow season). "
    
    for cp in calibration_points: 
        num_seasons = seasons.at[cp, 'NumSeasons']
        
        
        slide = prs.slides.add_slide(title_slide_layout)
        img = outdir +cp.replace(" ", "_") + ".png"
        pic = slide.shapes.add_picture(img, left= Inches(0.25), top= Inches(1.2), height= Inches(6.1))
                
    
        
        
        
        top = Inches(0.1)
        left = height = Inches(0.25)
        width= Inches (5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        text = cp.title() + "\nNumber of Seasons: " + str(num_seasons)
        if num_seasons==2: 
            text = text+ "\nWinter Selected Years: "+ str(events_df[cp + "_winter"].dropna().astype(int).tolist()) 
        
        text = text + "\nSpring Selected Years: " +   str(events_df[cp].dropna().astype(int).tolist()) 
    
        tf.text = text
            
            
     
    prs.save(outdir + 'EventSelectionTimeseriesPlots.pptx')
    
    
else: 
    list_seasons = ['winter', 'spring']
    for seas in list_seasons: 
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[6]
        print("Generating Summary PPTX")
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Selected Events Precip/Tmax/Flow"
        if seas == 'spring': 
            subtitle.text = "\nSpring Season Only: Plotted from Apr 1- Sep 30 (To show full flow season). "
        elif seas == 'winter':
            subtitle.text = "\nWinter Season Only: Plotted from Oct 1- Mar 30"
        
        for cp in calibration_points: 
            num_seasons = seasons.at[cp, 'NumSeasons']
            if num_seasons ==1 and seas == "winter" : continue
            img = outdir +cp.replace(" ", "_") + "_" + seas+ ".png"
            slide = prs.slides.add_slide(title_slide_layout)
            pic = slide.shapes.add_picture(img, left= Inches(0.25), top= Inches(1.2), height= Inches(6.1))
                    
        
            top = Inches(0.1)
            left = height = Inches(0.25)
            width= Inches (5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            text = cp.title() + "\nNumber of Seasons: " + str(num_seasons)
            if seas == "spring":
                text = text + "\nSpring Selected Years: " +   str(events_df[cp].dropna().astype(int).tolist()) 
            elif seas == "winter":
                test = text + "\nWinter Selected Years: " +   str(events_df[cp+ "_winter"].dropna().astype(int).tolist()) 
            tf.text = text
                
                
         
        prs.save(outdir + 'EventSelectionTimeseriesPlots_'+ seas.title() + 'Only.pptx')
    
    
    
    
#%% Compare nested vs indiv (Can only be run after both "nested" and "indiv" options above have been run) and write into a word document

from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from docx import Document
from docx.shared import Inches

# Function to change orientation of pg
#source: https://stackoverflow.com/questions/31893557/python-docx-sections-page-orientation
def change_orientation2():
    current_section = document.sections[-1]
    new_width, new_height = current_section.page_height, current_section.page_width
    # new_section = document.add_section(WD_SECTION.NEW_PAGE)
    current_section.orientation = WD_ORIENT.LANDSCAPE
    current_section.page_width = new_width
    current_section.page_height = new_height

    return current_section
   


events_df_nested = pd.read_csv(event_dir)
events_df=events_df.dropna(how='all')
events_df_indiv = pd.read_csv(event_indiv)
events_df_indiv= events_df_indiv.iloc[:10]

col_names = events_df_nested.columns.tolist()
col_names = [col.lower().replace("spaulding", "spalding") for col in col_names ]
events_df_nested.columns = col_names

col_names = events_df_indiv.columns.tolist()
col_names = [col.lower().replace("spaulding", "spalding") for col in col_names ]
events_df_indiv.columns = col_names


list_seasons = ['spring', 'winter']
doc_name = r"\\westfolsom\Projects\2021\CRT_Task3\Memos\EventSelection\Appendix\AppendixE.docx"
document = Document()
change_orientation2()
document.add_heading("Appendix E: Precipitation, Temperature, Flow, and SWE Variability for Selected Events")
i=0
for seas in list_seasons: 
    print(seas)

    for cp in calibration_points: 
        num_seasons = seasons.at[cp, 'NumSeasons']
        
        img_nested = outdir_nested +cp.replace(" ", "_") + "_" + seas+ ".png"
        img_indiv = outdir_indiv + cp.replace(" ", "_") + "_" + seas + ".png"
        
        if (os.path.exists(img_nested))&(os.path.exists(img_indiv)):
            i+=1
        # Join two plots into one
            images = [Image.open(x) for x in [img_indiv ,img_nested]]
            widths, heights = zip(*(i.size for i in images))
            
            total_width = sum(widths)
            max_height = max(heights)
            
            new_im = Image.new('RGB', (total_width, max_height))
            
            x_offset = 0
            for im in images:
              new_im.paste(im, (x_offset,0))
              x_offset += im.size[0]
            
            figname=outdir_both + cp.replace(" ", "_") + "_"+seas+".png"
            new_im.save(figname)
    
        
            # write to document
            p = document.add_paragraph()
            p.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
            r = p.add_run()
            document.add_picture(figname, width=Inches(9))
            
            caption_str = "Precipitation, temperature, flow, and SWE for individual (i 1 to 4) and nested (n 1 to 4) selected events for "+ str.capitalize(cp) + " for the " +seas + " season."
        
            p = document.add_paragraph("Figure E."+ str(i) + ": " + caption_str)
            
document.save(doc_name)
    


#%% Compare nested vs indiv (Can only be run after both "nested" and "indiv" options above have been run)

events_df_nested = pd.read_csv(event_dir)
events_df_indiv = pd.read_csv(event_indiv)
events_df_indiv= events_df_indiv.iloc[:10]

col_names = events_df_nested.columns.tolist()
col_names = [col.lower().replace("spaulding", "spalding") for col in col_names ]
events_df_nested.columns = col_names

col_names = events_df_indiv.columns.tolist()
col_names = [col.lower().replace("spaulding", "spalding") for col in col_names ]
events_df_indiv.columns = col_names


if plotTogether: 


    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]
    print("Generating Summary PPTX ")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Selected Events Precip/Tmax/Flow"
    subtitle.text = "Winter Season: Plotted from Oct 1- Mar 31\nSpring Season: Plotted from Apr 1- Sep 30 (To show full flow season). \nDate Lines accurate for NON-LEAP YEARS"
    
    #Loop through each calibration point
    for cp in calibration_points: 
        num_seasons = seasons.at[cp, 'NumSeasons']
        if num_seasons ==2: 
            top_in = 2
            height_in = 4
        else: 
            top_in = 1.5
            height_in = 4.8
        
        slide = prs.slides.add_slide(title_slide_layout)
        img_nested = outdir_nested +cp.replace(" ", "_") + ".png"
        img_indiv = outdir_indiv + cp.replace(" ", "_") + ".png"
        
        
        pic = slide.shapes.add_picture(img_indiv, left= Inches(0.25), top= Inches(top_in), height= Inches(height_in))
        pic = slide.shapes.add_picture(img_nested, left= Inches(5.1), top= Inches(top_in), height= Inches(height_in))
    
    
        
        
        
        top = Inches(0.1)
        left = height = Inches(0.25)
        width= Inches (5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        text = cp.title() + "\nNumber of Seasons: " + str(num_seasons)
        if num_seasons==2: 
            
            text = text+ "\nWinter Selected Years (Indiv): "+ str(events_df_indiv[cp + "_winter"].dropna().astype(int).tolist()) 
            text = text+ "\nWinter Selected Years (Nested): "+ str(events_df_nested[cp + "_winter"].dropna().astype(int).tolist()) 
        
        
        text = text + "\nSpring Selected Years (Indiv): " +   str(events_df_indiv[cp].dropna().astype(int).tolist()) 
        text = text + "\nSpring Selected Years (Nested): " +   str(events_df_nested[cp].dropna().astype(int).tolist()) 
    
        text = text + "\nIndividual\t\t\t\t Nested"
        tf.text = text
            
            
     
    prs.save(outdir_nested + 'NestedVSIndiv_EventSelectionTimeseriesPlots_updated.pptx')
    
    
else: 
    list_seasons = ['winter', 'spring']
    for seas in list_seasons: 
        
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[6]
        print("Generating Summary PPTX")
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Selected Events Precip/Tmax/Flows"
        if seas == 'spring': 
            subtitle.text = "\nSpring Season Only: Plotted from Apr 1- Sep 30 (To show full flow season). "
        elif seas == 'winter':
            subtitle.text = "\nWinter Season Only: Plotted from Oct 1- Mar 30"
        
        for cp in calibration_points: 
            num_seasons = seasons.at[cp, 'NumSeasons']
            if num_seasons ==1 and seas == "winter" : continue
            img = outdir +cp.replace(" ", "_") + "_" + seas+ ".png"
            slide = prs.slides.add_slide(title_slide_layout)
            
            
            img_nested = outdir_nested +cp.replace(" ", "_") + "_" + seas+ ".png"
            img_indiv = outdir_indiv + cp.replace(" ", "_") + "_" + seas + ".png"
        
            top_in = 2.2
            height_in = 4.8
            
            pic = slide.shapes.add_picture(img_indiv, left= Inches(0.25), top= Inches(top_in), height= Inches(height_in))
            pic = slide.shapes.add_picture(img_nested, left= Inches(5.1), top= Inches(top_in), height= Inches(height_in))
    
            # pic = slide.shapes.add_picture(img, left= Inches(0.25), top= Inches(1.2), height= Inches(6.1))
                    
        
            top = Inches(0.1)
            left = height = Inches(0.25)
            width= Inches (5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            text = cp.title() + "\nNumber of Seasons: " + str(num_seasons)

            if seas == 'winter': 
                
                text = text+ "\nWinter Selected Years (Indiv): "+ str(events_df_indiv[cp + "_winter"].dropna().astype(int).tolist()) 
                text = text+ "\nWinter Selected Years (Nested): "+ str(events_df_nested[cp + "_winter"].dropna().astype(int).tolist()) 
            
            elif seas == 'spring':
                text = text + "\nSpring Selected Years (Indiv): " +   str(events_df_indiv[cp].dropna().astype(int).tolist()) 
                text = text + "\nSpring Selected Years (Nested): " +   str(events_df_nested[cp].dropna().astype(int).tolist()) 
        
            text = text + "\n\n\t\t\tIndividual\t\t\t\t\t\t\t\t\tNested"
            tf.text = text
                
                
         
        prs.save(outdir + 'NestedVSIndivEventSelectionTimeseriesPlots_'+ seas.title() + 'Only.pptx')























